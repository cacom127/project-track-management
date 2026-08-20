"""CHANGE-017 (EXPORT-05/06/07) — test đơn vị cho `build_presentation`,
không cần DB/HTTP (chỉ cần list `ProjectOut` dựng tay). Dùng file
template thật (`app/export/assets/template.pptx`) — không mock nội
dung `python-pptx`, chỉ mock `app.core.s3.get_object_bytes` (không có
S3 emulator local, cùng tinh thần với test attachment hiện có)."""

import io
from datetime import date, datetime
from decimal import Decimal
from unittest.mock import MagicMock

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from app.export import service
from app.projects.schemas import ProjectOut


def _fake_png_bytes() -> bytes:
    # add_picture() dùng PIL để đọc kích thước ảnh gốc — cần bytes ảnh
    # hợp lệ thật (không phải b"x") để không lỗi UnidentifiedImageError.
    buffer = io.BytesIO()
    Image.new("RGB", (4, 4), color="red").save(buffer, format="PNG")
    return buffer.getvalue()


def _make_project(project_id: int, **overrides) -> ProjectOut:
    base = dict(
        id=project_id,
        customer_name="極秘商事株式会社",
        project_name=f"プロジェクト{project_id}",
        description="概要テキスト",
        start_date=date(2024, 1, 1),
        end_date=None,
        is_ongoing=True,
        team_size=5,
        total_man_month=Decimal("10.5"),
        source_note="確認元の内部メモ",
        created_by="user-1",
        created_at=datetime(2024, 1, 1),
        updated_at=datetime(2024, 1, 1),
        technologies=["React", "AWS"],
        project_types=["offshore", "ses"],
        industry="製造業",
        outcome_note="成果テキスト",
        dev_process_phases=["requirements", "design"],
        team_composition_note="チーム体制の詳細メモ",
    )
    base.update(overrides)
    return ProjectOut(**base)


def _all_text(slide) -> str:
    return "\n".join(shp.text_frame.text for shp in slide.shapes if shp.has_text_frame)


def test_one_slide_per_project(monkeypatch):
    monkeypatch.setattr(service.s3, "get_object_bytes", MagicMock(return_value=_fake_png_bytes()))
    projects = [_make_project(1), _make_project(2), _make_project(3)]

    data = service.build_presentation(projects, {})

    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3


def test_max_4_images_embedded_even_with_more_attachments(monkeypatch):
    fake_bytes = MagicMock(return_value=_fake_png_bytes())
    monkeypatch.setattr(service.s3, "get_object_bytes", fake_bytes)
    project = _make_project(1)
    s3_keys = [f"projects/1/{i}.jpg" for i in range(6)]

    data = service.build_presentation([project], {1: s3_keys})

    prs = Presentation(io.BytesIO(data))
    slide = prs.slides[0]
    pictures = [shp for shp in slide.shapes if shp.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 4
    assert fake_bytes.call_count == 4


def test_zero_attachments_renders_without_error(monkeypatch):
    monkeypatch.setattr(service.s3, "get_object_bytes", MagicMock())

    data = service.build_presentation([_make_project(1)], {})

    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 1


def test_excludes_customer_name_and_internal_notes(monkeypatch):
    monkeypatch.setattr(service.s3, "get_object_bytes", MagicMock())
    project = _make_project(1)

    data = service.build_presentation([project], {})

    prs = Presentation(io.BytesIO(data))
    text = _all_text(prs.slides[0])
    assert project.customer_name not in text
    assert project.source_note not in text
    assert project.team_composition_note not in text
    # sanity: field ĐƯỢC phép hiển thị vẫn có mặt (không phải test rỗng do lỗi)
    assert project.project_name in text
    assert project.outcome_note in text


def test_long_tech_badges_wrap_to_multiple_lines_without_overlapping(monkeypatch):
    # Feedback thực tế: badge 技術 dài đè lên badge kế tiếp vì code cũ
    # tính vị trí theo chiều rộng CỐ ĐỊNH của badge mẫu, không theo độ
    # dài text thật.
    monkeypatch.setattr(service.s3, "get_object_bytes", MagicMock())
    long_tags = [f"非常に長い技術タグ名前です{i}" for i in range(10)]
    project = _make_project(1, technologies=long_tags, dev_process_phases=[])

    data = service.build_presentation([project], {})

    prs = Presentation(io.BytesIO(data))
    slide = prs.slides[0]
    tech_badges = [shp for shp in slide.shapes if shp.name == "field_tech_badge_1"]
    assert len(tech_badges) == 10

    by_line: dict[int, list] = {}
    for badge in tech_badges:
        by_line.setdefault(badge.top, []).append(badge)
    assert len(by_line) > 1, "10 badge tên dài phải tự xuống ít nhất 2 dòng, không dồn 1 dòng"

    for line_badges in by_line.values():
        line_badges.sort(key=lambda b: b.left)
        for left_badge, right_badge in zip(line_badges, line_badges[1:]):
            assert left_badge.left + left_badge.width <= right_badge.left, (
                "2 badge cùng dòng không được chồng lên nhau theo chiều ngang"
            )


def test_many_long_badges_push_outcome_section_down_but_stays_within_slide(monkeypatch):
    monkeypatch.setattr(service.s3, "get_object_bytes", MagicMock())
    long_tags = [f"非常に長い技術タグ{i}" for i in range(8)]
    long_phases = [f"非常に長い開発工程名{i}" for i in range(6)]
    project = _make_project(1, technologies=long_tags, dev_process_phases=long_phases)

    data = service.build_presentation([project], {})

    prs = Presentation(io.BytesIO(data))
    slide = prs.slides[0]
    outcome = next(shp for shp in slide.shapes if shp.name == "field_outcome_note")
    label_outcome = next(shp for shp in slide.shapes if shp.name == "label_outcome")
    dividers = sorted((shp for shp in slide.shapes if shp.name == "divider"), key=lambda s: s.top)

    # divider cuối (trước 成果・課題・解決策) phải bị đẩy xuống thấp hơn
    # vị trí tĩnh gốc (5.65in) vì có nhiều badge dài phía trên.
    assert dividers[2].top > Emu(Inches(5.65))
    assert label_outcome.top > dividers[2].top
    assert outcome.top > label_outcome.top
    # vẫn nằm trong slide (7.5in), có chiều cao tối thiểu hợp lý.
    assert outcome.top + outcome.height <= service.SLIDE_HEIGHT_EMU
    assert outcome.height >= service.MIN_OUTCOME_HEIGHT_EMU


def test_badge_count_matches_actual_values_not_template_sample_count(monkeypatch):
    monkeypatch.setattr(service.s3, "get_object_bytes", MagicMock())
    project = _make_project(
        1,
        technologies=["React", "AWS", "Java", "Python", "Go"],
        project_types=[],
        dev_process_phases=["requirements"],
    )

    data = service.build_presentation([project], {})

    prs = Presentation(io.BytesIO(data))
    slide = prs.slides[0]
    tech_badges = [shp for shp in slide.shapes if shp.name == "field_tech_badge_1"]
    phase_badges = [shp for shp in slide.shapes if shp.name == "field_phase_badge_1"]
    type_badges = [shp for shp in slide.shapes if shp.name == "field_type_badge_1"]
    assert len(tech_badges) == 5
    assert len(phase_badges) == 1
    assert len(type_badges) == 0
    status_badges = [shp for shp in slide.shapes if shp.name == "field_status_badge"]
    assert len(status_badges) == 1
    assert status_badges[0].text_frame.text == "進行中"
