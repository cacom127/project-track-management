"""CHANGE-017 (EXPORT-05..08) — dựng file `.pptx` từ template đã thiết
kế sẵn (`assets/template.pptx`, xem `changes/CHANGE-017-project-export-
pptx/tasks.md` T2 cho danh sách tên shape).

Cách "nhân bản slide": python-pptx không có `duplicate_slide()` có sẵn.
Cách chuẩn (và cách dùng ở đây): `add_slide(layout)` để tạo 1 slide rỗng
cùng layout, rồi `copy.deepcopy` từng shape-element của slide mẫu và
append vào `<p:spTree>` của slide mới — layout `template.pptx` dùng là
"Blank" (không có placeholder tự sinh) nên slide mới tạo ra hoàn toàn
trống trước khi copy, không lo xung đột placeholder.

Mỗi project được apply lên 1 slide RIÊNG (nhân bản từ slide mẫu gốc,
CHƯA bị điền dữ liệu) — slide mẫu gốc bị xoá ở cuối cùng, tránh phụ
thuộc vào thứ tự điền dữ liệu (không có project nào "dùng chung" slide
mẫu ban đầu).

Layout CASCADING (feedback thực tế: badge 技術/開発工程 dài đè lên
badge khác vì code cũ tính vị trí badge kế tiếp theo CHIỀU RỘNG CỐ ĐỊNH
của badge mẫu, không theo độ dài text thật):
  - Mọi hàng badge (種別+trạng thái ở header, 技術, 開発工程) đều tính
    chiều rộng từng badge theo số ký tự thật (`_estimate_badge_width`)
    và TỰ XUỐNG DÒNG nếu vượt giới hạn phải của hàng (`_flow_place`).
  - Vì số dòng badge thay đổi theo từng project, các phần tử tĩnh phía
    dưới (2 divider còn lại, khối 業種/期間/..., cột ảnh, label các
    hàng badge, label/khung 成果・課題・解決策) đều được ĐẶT LẠI vị trí
    (top) theo kiểu "cursor" tuần tự từ trên xuống, KHÔNG dùng toạ độ
    tĩnh của template nữa — xem `_fill_slide`.
  - `field_outcome_note` height co lại tương ứng nếu phần trên chiếm
    nhiều chỗ hơn bình thường, đảm bảo luôn nằm trong slide (trừ
    trường hợp cực đoan — clamp tối thiểu `MIN_OUTCOME_HEIGHT_EMU`).

Đồng thời sửa 1 bug liên quan phát hiện khi rà lại: `field_project_name`/
`field_meta`/các label vốn mang `auto_size=SHAPE_TO_FIT_TEXT` (mặc định
của `python-pptx` khi tạo text box, KHÔNG phải cố ý) — nghĩa là khung
tự PHÌNH RA theo text dài thay vì co chữ lại, có thể đè xuống phần tử
bên dưới. Đã patch trực tiếp trong `assets/template.pptx`: `field_meta`
đổi sang co chữ (`TEXT_TO_FIT_SHAPE`, giống `field_description`/
`field_outcome_note`); `field_project_name` và label đổi sang `NONE`
(cố định, không phình không co — giữ đúng yêu cầu "cỡ chữ CỐ ĐỊNH"
riêng cho title, xem EXPORT-08)."""

import copy
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

from app.core import s3
from app.export.labels import DEV_PROCESS_PHASE_LABELS, PROJECT_TYPE_LABELS
from app.projects.schemas import ProjectOut

TEMPLATE_PATH = Path(__file__).parent / "assets" / "template.pptx"

MAX_IMAGES_PER_SLIDE = 4
SLIDE_HEIGHT_EMU = Inches(7.5)

# --- Hằng số layout cascading (đơn vị EMU, 914400 EMU = 1 inch) ---
HEADER_BADGES_TOP_EMU = Inches(1.05)  # cố định, ngay dưới title (title KHÔNG phình — xem docstring)
HEADER_LEFT_EMU = Inches(0.4)  # trái của title/badge 種別+trạng thái
ROW_RIGHT_LIMIT_EMU = Inches(12.93)  # khớp mép phải divider (0.4 + 12.53)
# 技術/開発工程 nằm NGANG NHAU thành 2 cột (feedback thực tế — trước đó
# xếp trên-dưới) — mỗi cột rộng bằng nhau, có khoảng hở giữa 2 cột.
_BADGE_ROW_LEFT_EMU = Inches(0.4)
BADGE_COL_GAP_EMU = Inches(0.3)
_BADGE_COL_WIDTH_EMU = (ROW_RIGHT_LIMIT_EMU - _BADGE_ROW_LEFT_EMU - BADGE_COL_GAP_EMU) // 2
BADGE_COL_LEFT_EMU = [
    _BADGE_ROW_LEFT_EMU,
    _BADGE_ROW_LEFT_EMU + _BADGE_COL_WIDTH_EMU + BADGE_COL_GAP_EMU,
]
BADGE_COL_RIGHT_LIMIT_EMU = [
    BADGE_COL_LEFT_EMU[0] + _BADGE_COL_WIDTH_EMU,
    ROW_RIGHT_LIMIT_EMU,
]
# Từ label-left tới badge-left trong CÙNG 1 cột — cũng là mức thụt
# dòng khi wrap.
BADGE_INDENT_EMU = Inches(1.0)
BADGE_HEIGHT_EMU = Inches(0.34)
BADGE_LINE_GAP_EMU = Inches(0.10)  # khoảng cách dọc giữa 2 dòng badge khi wrap
BADGE_GAP_EMU = Inches(0.12)  # khoảng cách ngang giữa 2 badge cùng dòng
BADGE_H_PADDING_EMU = Inches(0.16)  # padding 2 bên trong 1 badge, cộng theo text
# Sàn tối thiểu CHUNG cho mọi badge — CỐ Ý nhỏ (không dùng chiều rộng
# badge mẫu trong template làm sàn như trước) để badge LUÔN phản ánh
# đúng độ dài text thật. Bug thực tế: catalog 種別 (`オフショア`/`SES`/
# `ラボ`/`新規開発`/`保守`) toàn nhãn ngắn, sàn cũ (theo template,
# ~1.6in) luôn thắng ước lượng theo text → nhìn như badge 種別 không
# đổi theo độ dài, dù công thức vẫn tính đúng bên dưới.
BADGE_MIN_WIDTH_EMU = Inches(0.5)
# Ước lượng rộng/ký tự — dư cho CJK+Latin trộn (thà badge hơi rộng còn
# hơn chữ bị cắt).
CHAR_WIDTH_EMU = Inches(0.11)
# Khoảng cách giữa 2 khối lớn (2 cột -> badge, badge -> badge, badge -> divider).
GROUP_GAP_EMU = Inches(0.15)
DIVIDER_CONTENT_GAP_EMU = Inches(0.09)  # khoảng cách từ divider tới nội dung ngay sau nó
DESC_HEIGHT_EMU = Inches(1.55)
META_GAP_EMU = Inches(0.15)  # khoảng cách từ dưới field_description tới field_meta
META_HEIGHT_EMU = Inches(1.0)
IMG_HEIGHT_EMU = Inches(1.5)
IMG_ROW_GAP_EMU = Inches(0.12)
OUTCOME_LABEL_GAP_EMU = Inches(0.05)
BOTTOM_MARGIN_EMU = Inches(0.3)
MIN_OUTCOME_HEIGHT_EMU = Inches(0.5)


def _find_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _find_shapes(slide, name: str) -> list:
    return [shape for shape in slide.shapes if shape.name == name]


def _remove_shape(shape) -> None:
    if shape is not None:
        shape._element.getparent().remove(shape._element)


def _set_text(shape, text: str) -> None:
    """Set text mà vẫn giữ style (font/size/color/bold) của run đầu
    tiên đã thiết kế sẵn trong template — chỉ đổi nội dung."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
        run.text = text
    else:
        run = p.add_run()
        run.text = text


def _duplicate_shape(after_shape):
    """Deepcopy XML element của `after_shape`, append ngay sau nó trong
    cùng slide, trả về Shape wrapper của bản copy. Vị trí/kích thước
    của bản copy do caller set sau (`_flow_place`) — không set ở đây.

    QUAN TRỌNG: khi cần nhân bản NHIỀU lần cho cùng 1 tên shape (N
    badge), phải gọi nối chuỗi — `_duplicate_shape(clone_truoc)`, KHÔNG
    lặp lại `_duplicate_shape(master)` — nếu không, `addnext` luôn chèn
    ngay sau `master` nên các lần gọi sau sẽ chèn NGƯỢC thứ tự, khiến
    việc tìm "bản clone vừa tạo" (`matches[-1]`) lấy nhầm shape cũ, bỏ
    sót 1 clone không bao giờ được set lại vị trí (bug thật đã phát
    hiện qua test `test_long_tech_badges_wrap_to_multiple_lines_...`)."""
    clone_el = copy.deepcopy(after_shape._element)
    after_shape._element.addnext(clone_el)
    slide_shapes = after_shape._parent
    matches = [shp for shp in slide_shapes if shp.name == after_shape.name]
    return matches[-1]


def _estimate_badge_width(text: str) -> int:
    """Ước lượng chiều rộng badge theo SỐ KÝ TỰ THẬT — thay cho việc
    dùng cố định chiều rộng badge mẫu trong template (nguyên nhân badge
    dài đè lên badge kế tiếp, feedback thực tế). Sàn tối thiểu CỐ ĐỊNH
    nhỏ (`BADGE_MIN_WIDTH_EMU`), KHÔNG dùng chiều rộng badge mẫu làm
    sàn — nếu không, badge có nhãn ngắn (vd 種別) sẽ luôn bị sàn theo
    template che mất, trông như không đổi theo độ dài text."""
    return max(BADGE_MIN_WIDTH_EMU, BADGE_H_PADDING_EMU + len(text) * CHAR_WIDTH_EMU)


def _flow_place(
    shapes_with_widths: list[tuple], start_left: int, start_top: int, right_limit: int
) -> int:
    """Đặt tuần tự các shape (đã có text) theo chiều ngang, tự xuống
    dòng khi vượt `right_limit` — KHÔNG cắt/bỏ shape nào (khác với
    hướng "+N" đã cân nhắc, feedback thực tế chọn giữ đủ thông tin).
    Trả về top NGAY SAU khối này, để phần tử kế tiếp không đè lên dù số
    dòng thay đổi theo từng project."""
    if not shapes_with_widths:
        return start_top

    x = start_left
    line = 0
    for shape, width in shapes_with_widths:
        if x != start_left and x + width > right_limit:
            line += 1
            x = start_left
        y = start_top + line * (BADGE_HEIGHT_EMU + BADGE_LINE_GAP_EMU)
        shape.left = Emu(x)
        shape.top = Emu(y)
        shape.width = Emu(width)
        x += width + BADGE_GAP_EMU

    total_lines = line + 1
    return start_top + total_lines * BADGE_HEIGHT_EMU + (total_lines - 1) * BADGE_LINE_GAP_EMU


def _fill_header_badges(slide, project: ProjectOut, top: int) -> int:
    """種別 (N badge, N thay đổi theo project) + trạng thái (1 badge cố
    định) — CÙNG 1 luồng wrap, không có label đứng trước (khác hàng
    技術/開発工程)."""
    type_master = _find_shape(slide, "field_type_badge_1")
    status_shape = _find_shape(slide, "field_status_badge")

    type_labels = [PROJECT_TYPE_LABELS.get(code, code) for code in project.project_types]
    items: list[tuple] = []
    if type_labels:
        last = type_master
        for index, label in enumerate(type_labels):
            shape = type_master if index == 0 else _duplicate_shape(last)
            last = shape
            _set_text(shape, label)
            items.append((shape, _estimate_badge_width(label)))
    else:
        _remove_shape(type_master)

    status_label = "進行中" if project.is_ongoing else "終了"
    _set_text(status_shape, status_label)
    items.append((status_shape, _estimate_badge_width(status_label)))

    return _flow_place(items, HEADER_LEFT_EMU, top, ROW_RIGHT_LIMIT_EMU)


def _fill_badge_row(
    slide,
    label_name: str,
    master_name: str,
    extra_name: str,
    values: list[str],
    top: int,
    *,
    col_left: int,
    right_limit: int,
) -> int:
    """1 cột badge có label đứng trước (技術:/開発工程:) — tự xuống
    dòng TRONG PHẠM VI CỘT (`col_left`..`right_limit`), dòng wrap thụt
    vào ngang mức badge đầu (không thụt theo label). Trả về top NGAY
    SAU cột này (để caller tính `max()` giữa 2 cột — số dòng của 技術 và
    開発工程 độc lập nhau, feedback thực tế: đặt 2 cột NGANG NHAU thay
    vì trên-dưới)."""
    label = _find_shape(slide, label_name)
    master = _find_shape(slide, master_name)
    extra = _find_shape(slide, extra_name)
    # Template chỉ có 2 badge mẫu (field_*_badge_1/2) — không cần mẫu
    # thứ 2 nữa vì số badge thật được clone động từ mẫu 1 (PROJ-30 style).
    _remove_shape(extra)

    if not values:
        _remove_shape(label)
        _remove_shape(master)
        return top

    label.left = Emu(col_left)
    label.top = Emu(top)

    items = []
    last = master
    for index, value in enumerate(values):
        shape = master if index == 0 else _duplicate_shape(last)
        last = shape
        _set_text(shape, value)
        items.append((shape, _estimate_badge_width(value)))

    badge_left = col_left + BADGE_INDENT_EMU
    return _flow_place(items, badge_left, top, right_limit)


def _format_period(project: ProjectOut) -> str:
    start = project.start_date.isoformat()
    if project.is_ongoing:
        return f"{start} 〜 進行中"
    if project.end_date:
        return f"{start} 〜 {project.end_date.isoformat()}"
    return start


def _fill_meta(slide, project: ProjectOut, top: int) -> None:
    shape = _find_shape(slide, "field_meta")
    shape.top = Emu(top)
    industry = project.industry or "—"
    period = _format_period(project)
    team_size = f"{project.team_size}名" if project.team_size is not None else "—"
    man_month = f"{project.total_man_month}人月" if project.total_man_month is not None else "—"
    _set_text(
        shape,
        f"業種: {industry}　　期間: {period}\n人数: {team_size}　　総人月: {man_month}",
    )


def _fill_images(slide, s3_keys: list[str], top: int) -> None:
    # EXPORT-07 — tối đa 4 ảnh đầu; slot thừa (project ít ảnh hơn 4) giữ
    # nguyên placeholder xám của template, không lỗi.
    row_tops = [top, top + IMG_HEIGHT_EMU + IMG_ROW_GAP_EMU]
    for index in range(MAX_IMAGES_PER_SLIDE):
        slot_name = f"img_slot_{index + 1}"
        slot = _find_shape(slide, slot_name)
        if slot is None:
            continue
        slot.top = Emu(row_tops[index // 2])
        if index >= len(s3_keys):
            continue
        left, slot_top, width, height = slot.left, slot.top, slot.width, slot.height
        image_bytes = s3.get_object_bytes(s3_keys[index])
        _remove_shape(slot)
        slide.shapes.add_picture(io.BytesIO(image_bytes), left, slot_top, width, height)


def _fill_slide(slide, project: ProjectOut, s3_keys: list[str]) -> None:
    # 3 divider tĩnh trong template, LẤY THỨ TỰ theo vị trí gốc (top)
    # TRƯỚC khi bất kỳ shape nào bị di chuyển — dividers[1]/[2] sẽ được
    # đặt lại top ở dưới theo nội dung thực tế của slide này.
    dividers = sorted(_find_shapes(slide, "divider"), key=lambda shape: shape.top)
    divider_after_header, divider_before_badges, divider_before_outcome = dividers

    _set_text(_find_shape(slide, "field_project_name"), project.project_name)
    header_bottom = _fill_header_badges(slide, project, HEADER_BADGES_TOP_EMU)

    divider_after_header.top = Emu(header_bottom + GROUP_GAP_EMU)
    two_col_top = divider_after_header.top + DIVIDER_CONTENT_GAP_EMU

    desc_shape = _find_shape(slide, "field_description")
    desc_shape.top = Emu(two_col_top)
    _set_text(desc_shape, project.description or "—")
    _fill_meta(slide, project, two_col_top + DESC_HEIGHT_EMU + META_GAP_EMU)
    _fill_images(slide, s3_keys[:MAX_IMAGES_PER_SLIDE], two_col_top)

    left_col_bottom = two_col_top + DESC_HEIGHT_EMU + META_GAP_EMU + META_HEIGHT_EMU
    right_col_bottom = two_col_top + 2 * IMG_HEIGHT_EMU + IMG_ROW_GAP_EMU
    # Bug thật phát hiện qua rà toạ độ: divider này trước đó KHÔNG được
    # đặt lại vị trí (chỉ 2/3 divider được cập nhật) — vẫn nằm ở toạ độ
    # tĩnh của template, có thể cắt ngang qua hàng badge khi nội dung
    # cột trên (概要/ảnh) cao/thấp hơn bình thường.
    divider_before_badges.top = Emu(max(left_col_bottom, right_col_bottom) + GROUP_GAP_EMU)
    badges_top = divider_before_badges.top + DIVIDER_CONTENT_GAP_EMU

    # 技術/開発工程 nằm NGANG NHAU thành 2 cột (feedback thực tế) — mỗi
    # cột tự xuống dòng độc lập, dùng max() 2 đáy để tính phần tử kế
    # tiếp, KHÔNG cộng dồn tuần tự như 2 cột 概要/ảnh phía trên.
    tech_bottom = _fill_badge_row(
        slide,
        "label_tech",
        "field_tech_badge_1",
        "field_tech_badge_2",
        list(project.technologies),
        badges_top,
        col_left=BADGE_COL_LEFT_EMU[0],
        right_limit=BADGE_COL_RIGHT_LIMIT_EMU[0],
    )
    phase_labels = [DEV_PROCESS_PHASE_LABELS.get(code, code) for code in project.dev_process_phases]
    phase_bottom = _fill_badge_row(
        slide,
        "label_phase",
        "field_phase_badge_1",
        "field_phase_badge_2",
        phase_labels,
        badges_top,
        col_left=BADGE_COL_LEFT_EMU[1],
        right_limit=BADGE_COL_RIGHT_LIMIT_EMU[1],
    )

    divider_before_outcome.top = Emu(max(tech_bottom, phase_bottom) + GROUP_GAP_EMU)
    label_outcome = _find_shape(slide, "label_outcome")
    label_outcome.top = divider_before_outcome.top + DIVIDER_CONTENT_GAP_EMU

    outcome_shape = _find_shape(slide, "field_outcome_note")
    outcome_top = label_outcome.top + label_outcome.height + OUTCOME_LABEL_GAP_EMU
    outcome_shape.top = Emu(outcome_top)
    outcome_shape.height = Emu(
        max(MIN_OUTCOME_HEIGHT_EMU, SLIDE_HEIGHT_EMU - BOTTOM_MARGIN_EMU - outcome_top)
    )
    _set_text(outcome_shape, project.outcome_note or "—")
    # EXPORT-06 — customer_name/source_note/team_composition_note KHÔNG
    # có shape tương ứng trong template nên không cần code loại trừ
    # riêng, chỉ cần không bao giờ đọc field này ở trên.


def _duplicate_slide(prs: Presentation, source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))
    return new_slide


def _delete_slide(prs: Presentation, slide) -> None:
    slide_id = slide.slide_id
    xml_slides = prs.slides._sldIdLst
    for sld_id_el in list(xml_slides):
        if int(sld_id_el.get("id")) == slide_id:
            r_id = sld_id_el.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            prs.part.drop_rel(r_id)
            xml_slides.remove(sld_id_el)
            return


def build_presentation(
    projects: list[ProjectOut], attachments_by_project: dict[int, list[str]]
) -> bytes:
    """EXPORT-05 — 1 slide/project. `attachments_by_project` map
    `project_id -> list[s3_key]` (đã giới hạn thứ tự upload, chưa cắt
    4 — `_fill_images` tự cắt)."""
    prs = Presentation(str(TEMPLATE_PATH))
    template_slide = prs.slides[0]

    for project in projects:
        slide = _duplicate_slide(prs, template_slide)
        _fill_slide(slide, project, attachments_by_project.get(project.id, []))

    _delete_slide(prs, template_slide)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
